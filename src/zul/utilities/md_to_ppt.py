import re
from pathlib import Path
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml import parse_xml
from pptx.util import Pt as PtUnit


class DynamicMarkdownToPPTXService:
    """
    Service dinamis untuk convert Markdown ke PPTX
    - Auto detect template (pakai template jika ada, generate baru jika tidak)
    - Customizable font styling per header level
    - Support image dari file path (no matplotlib execution)
    - Support table parsing dari markdown
    - Support convert to bytes untuk download
    - Support nested bullets dan section headers
    """

    def __init__(self, template_path=None, style_config=None):
        """
        Initialize service

        Args:
            template_path: Path ke template .pptx (optional)
            style_config: Dict untuk custom font styling (optional)
        """
        self.template_path = template_path
        self.use_template = False

        # Check template
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
            self.use_template = True
            print(f"✅ Mode: Template-based (using {template_path})")
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(5.625)
            print(f"✅ Mode: Generate from scratch")
            if template_path:
                print(f"   ⚠️  Template not found: {template_path}")

        # Set style configuration dengan merge
        self.style_config = self.get_default_style()
        if style_config:
            self._merge_style(style_config)

    def get_default_style(self):
        """Default styling configuration"""
        return {
            "h1": {
                "font_size": 44,
                "font_name": "Calibri",
                "font_color": RGBColor(0, 51, 102),
                "bold": True,
                "italic": False,
            },
            "h2": {
                "font_size": 32,
                "font_name": "Calibri",
                "font_color": RGBColor(230, 126, 34),
                "bold": True,
                "italic": False,
            },
            "h3": {
                "font_size": 24,
                "font_name": "Calibri",
                "font_color": RGBColor(52, 73, 94),
                "bold": True,
                "italic": False,
            },
            "h4": {
                "font_size": 20,
                "font_name": "Calibri",
                "font_color": RGBColor(52, 73, 94),
                "bold": True,
                "italic": False,
            },
            "body": {
                "font_size": 18,
                "font_name": "Calibri",
                "font_color": RGBColor(60, 60, 60),
                "bold": False,
                "italic": False,
            },
            "bullet": {
                "font_size": 18,
                "font_name": "Calibri",
                "font_color": RGBColor(60, 60, 60),
                "bold": False,
                "italic": False,
            },
            "table": {
                "font_size": 12,
                "font_name": "Calibri",
                "header_color": RGBColor(0, 51, 102),
                "cell_color": RGBColor(60, 60, 60),
                "bold_header": True,
            },
        }

    def _merge_style(self, new_style):
        """Deep merge new style dengan existing style"""
        for key, value in new_style.items():
            if key in self.style_config:
                self.style_config[key].update(value)
            else:
                self.style_config[key] = value

    def set_style(self, style_dict):
        """
        Update style configuration

        Example:
            service.set_style({
                'h1': {'font_size': 48, 'font_color': RGBColor(255, 0, 0)}
            })
        """
        self._merge_style(style_dict)
        print(f"✅ Style updated for: {', '.join(style_dict.keys())}")

    def apply_text_style(self, paragraph, style_type="body"):
        """Apply styling ke paragraph berdasarkan style type"""
        if style_type in self.style_config:
            style = self.style_config[style_type]
        else:
            style = self.style_config.get("body", {})

        if "font_size" in style and style["font_size"]:
            paragraph.font.size = Pt(style["font_size"])
        if "font_name" in style and style["font_name"]:
            paragraph.font.name = style["font_name"]
        if "font_color" in style and style["font_color"]:
            paragraph.font.color.rgb = style["font_color"]
        if "bold" in style:
            paragraph.font.bold = style["bold"]
        if "italic" in style:
            paragraph.font.italic = style["italic"]

    def enable_bullet(self, paragraph, level=0):
        """
        Enable bullet formatting untuk paragraph

        Args:
            paragraph: Paragraph object
            level: Bullet level (0, 1, 2, ...)
        """
        # Set paragraph level
        paragraph.level = level

        # Get or create pPr (paragraph properties)
        pPr = paragraph._element.get_or_add_pPr()

        # Remove any existing buNone (no bullet) element
        for buNone in pPr.findall(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"
        ):
            pPr.remove(buNone)

        # Add bullet char element if not exists
        buChar = pPr.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar"
        )
        if buChar is None:
            # Create bullet with default bullet character
            buChar_xml = '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>'
            buChar = parse_xml(buChar_xml)
            pPr.append(buChar)

    def create_blank_slide(self):
        """Create blank slide"""
        if self.use_template:
            try:
                blank_layout = self.prs.slide_layouts[6]
            except IndexError:
                blank_layout = self.prs.slide_layouts[-1]
            slide = self.prs.slides.add_slide(blank_layout)
        else:
            blank_layout = self.prs.slide_layouts[6]
            slide = self.prs.slides.add_slide(blank_layout)

            # Remove all shapes
            for shape in list(slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)

            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

        return slide

    def add_styled_textbox(
        self, slide, text, left, top, width, height, style_type="body"
    ):
        """Add textbox dengan custom styling"""
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )

        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.clear()

        p = text_frame.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.LEFT

        self.apply_text_style(p, style_type)

        return textbox

    def add_content_textbox(self, slide, content_items, left, top, width):
        """
        Add textbox dengan mixed content (sections, bullets, text)

        Args:
            slide: Slide object
            content_items: List of dicts with 'type' and 'text' and optional 'level'
            left, top: Position dalam inches
            width: Width dalam inches

        Returns:
            Total height used
        """
        # Estimate height
        total_height = len(content_items) * 0.45 + 0.5

        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(total_height)
        )

        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Clear text dari default paragraph
        if text_frame.paragraphs:
            text_frame.paragraphs[0].text = ""

        for idx, item in enumerate(content_items):
            # Use first paragraph for first item, add new for rest
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = item["text"]

            item_type = item.get("type", "bullet")
            level = item.get("level", 0)

            if item_type == "section":
                # Section header (### atau ####)
                p.level = 0
                style_type = "h3" if item.get("header_level", 3) == 3 else "h4"
                self.apply_text_style(p, style_type)
                p.space_before = Pt(12)
                p.space_after = Pt(6)
            elif item_type == "bullet":
                # CRITICAL: Enable bullet formatting explicitly
                self.enable_bullet(p, level)
                self.apply_text_style(p, "bullet")
            else:
                # Regular text
                p.level = 0
                self.apply_text_style(p, "body")

        return total_height

    def add_bullet_textbox(self, slide, bullets, left, top, width, style_type="bullet"):
        """Add textbox dengan bullet points (legacy support)"""
        total_height = len(bullets) * 0.4 + 0.2

        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(total_height)
        )

        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Clear text dari default paragraph
        if text_frame.paragraphs:
            text_frame.paragraphs[0].text = ""

        for idx, bullet_text in enumerate(bullets):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet_text

            # CRITICAL: Enable bullet formatting explicitly
            self.enable_bullet(p, level=0)
            self.apply_text_style(p, style_type)

        return total_height

    def parse_markdown_table(self, table_text):
        """
        Parse markdown table menjadi list of lists

        Args:
            table_text: String berisi markdown table

        Returns:
            List of lists [[header1, header2], [row1col1, row1col2], ...]
        """
        lines = [
            line.strip() for line in table_text.strip().split("\n") if line.strip()
        ]

        if len(lines) < 2:
            return None

        table_data = []

        for i, line in enumerate(lines):
            # Skip separator line (yang berisi |---|---|)
            if "---" in line:
                continue

            # Parse cells
            cells = [cell.strip() for cell in line.split("|") if cell.strip()]

            # Clean markdown formatting (bold, italic, etc)
            cleaned_cells = []
            for cell in cells:
                # Remove markdown bold (**text**)
                cell = re.sub(r"\*\*(.*?)\*\*", r"\1", cell)
                # Remove markdown italic (*text*)
                cell = re.sub(r"\*(.*?)\*", r"\1", cell)
                # Remove HTML breaks
                cell = cell.replace("<br>", "\n").replace("<br/>", "\n")
                cleaned_cells.append(cell.strip())

            table_data.append(cleaned_cells)

        return table_data if table_data else None

    def add_table(self, slide, table_data, left=0.5, top=1.5, width=9, height=None):
        """
        Add table ke slide

        Args:
            slide: Slide object
            table_data: List of lists [[header], [row1], [row2], ...]
            left, top: Position dalam inches
            width: Total width dalam inches
            height: Total height dalam inches (auto jika None)
        """
        if not table_data or len(table_data) == 0:
            return None

        rows = len(table_data)
        cols = len(table_data[0])

        # Calculate height
        if height is None:
            row_height = 0.4
            height = rows * row_height
        else:
            row_height = height / rows

        # Add table shape
        table_shape = slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        table = table_shape.table
        style = self.style_config.get("table", {})

        # Fill table data
        for i, row_data in enumerate(table_data):
            for j, cell_text in enumerate(row_data):
                if j < cols:
                    cell = table.cell(i, j)
                    cell.text = str(cell_text)

                    # Style cell text
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.LEFT

                        # Header row styling
                        if i == 0:
                            paragraph.font.bold = style.get("bold_header", True)
                            paragraph.font.size = Pt(style.get("font_size", 12))
                            paragraph.font.color.rgb = style.get(
                                "header_color", RGBColor(0, 51, 102)
                            )

                            # Header background color
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(220, 230, 241)
                        else:
                            paragraph.font.size = Pt(style.get("font_size", 12))
                            paragraph.font.color.rgb = style.get(
                                "cell_color", RGBColor(60, 60, 60)
                            )

                        if style.get("font_name"):
                            paragraph.font.name = style["font_name"]

        return table_shape

    def add_image(self, slide, image_path, left=1.5, top=3, width=7, height=None):
        """
        Add image ke slide dari file path

        Args:
            slide: Slide object
            image_path: Path ke image file
            left, top: Position dalam inches
            width: Width dalam inches
            height: Height dalam inches (optional, auto jika None)

        Returns:
            Picture shape atau None jika gagal
        """
        if not Path(image_path).exists():
            print(f"⚠️  Image not found: {image_path}")
            return None

        try:
            if height:
                pic = slide.shapes.add_picture(
                    image_path,
                    Inches(left),
                    Inches(top),
                    width=Inches(width),
                    height=Inches(height),
                )
            else:
                pic = slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top), width=Inches(width)
                )
            return pic
        except Exception as e:
            print(f"⚠️  Error adding image: {e}")
            return None

    def fill_template_placeholder(self, slide, title=None, content=None):
        """Fill placeholder di template slide"""
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type

            # Title placeholder
            if ph_type == 1 and title:
                shape.text = title
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text:
                            self.apply_text_style(paragraph, "h2")

            # Body/Content placeholder
            elif ph_type in [2, 7] and content:
                if not shape.has_text_frame:
                    continue

                text_frame = shape.text_frame
                text_frame.clear()

                if isinstance(content, list):
                    for idx, item in enumerate(content):
                        p = text_frame.add_paragraph()

                        # Check if item is dict (structured content)
                        if isinstance(item, dict):
                            p.text = item["text"]
                            item_type = item.get("type", "bullet")
                            level = item.get("level", 0)

                            if item_type == "section":
                                p.level = 0
                                style_type = (
                                    "h3" if item.get("header_level", 3) == 3 else "h4"
                                )
                                self.apply_text_style(p, style_type)
                            elif item_type == "bullet":
                                # Enable bullet for structured content
                                self.enable_bullet(p, level)
                                self.apply_text_style(p, "bullet")
                            else:
                                p.level = 0
                                self.apply_text_style(p, "body")
                        else:
                            p.text = str(item)
                            # Default to bullet for list items
                            self.enable_bullet(p, level=0)
                            self.apply_text_style(p, "bullet")
                else:
                    p = text_frame.add_paragraph()
                    p.text = str(content)
                    self.apply_text_style(p, "body")

    def add_slide_from_content(
        self,
        title=None,
        content=None,
        image_paths=None,
        table_data=None,
        layout_index=1,
    ):
        """
        Add slide dengan content

        Args:
            title: Slide title
            content: String, list, atau list of dicts untuk content
            image_paths: List of image paths atau single path
            table_data: List of lists untuk table data
            layout_index: Layout index jika pakai template
        """
        if self.use_template:
            # Mode template: gunakan layout dan fill placeholder
            try:
                layout = self.prs.slide_layouts[layout_index]
            except IndexError:
                layout = self.prs.slide_layouts[1]

            slide = self.prs.slides.add_slide(layout)
            self.fill_template_placeholder(slide, title, content)

        else:
            # Mode generate: buat manual dengan positioning
            slide = self.create_blank_slide()

            current_y = 0.5

            # Add title
            if title:
                self.add_styled_textbox(slide, title, 0.5, current_y, 9, 0.8, "h2")
                current_y += 1

            # Add content
            if content:
                if isinstance(content, list) and len(content) > 0:
                    # Check if structured content (list of dicts)
                    if isinstance(content[0], dict):
                        height = self.add_content_textbox(
                            slide, content, 0.8, current_y, 8.5
                        )
                    else:
                        height = self.add_bullet_textbox(
                            slide, content, 0.8, current_y, 8.5, "bullet"
                        )
                    current_y += height + 0.3
                elif isinstance(content, str):
                    self.add_styled_textbox(
                        slide, content, 0.5, current_y, 9, 0.6, "body"
                    )
                    current_y += 0.8

        # Add table
        if table_data:
            table_shape = self.add_table(
                slide,
                table_data,
                left=0.5,
                top=current_y if not self.use_template else 1.5,  # type: ignore
                width=9,
            )
            if table_shape and not self.use_template:
                current_y += table_shape.height.inches + 0.3  # type: ignore

        # Add images
        if image_paths:
            paths = image_paths if isinstance(image_paths, list) else [image_paths]

            for img_path in paths:
                if img_path:
                    pic = self.add_image(
                        slide,
                        img_path,
                        left=1.5,
                        top=3 if self.use_template else current_y,  # type: ignore
                        width=7,
                    )
                    if pic and not self.use_template:
                        current_y += pic.height.inches + 0.2  # type: ignore

        return slide

    def convert_markdown(self, md_content, output_path):
        """
        Convert markdown ke PPTX

        Supported syntax:
        - # Heading 1 (title slide)
        - ## Heading 2 (slide title)
        - ### Heading 3 (section header dalam slide)
        - #### Heading 4 (sub-section header dalam slide)
        - * Bullet point (level 0)
        - ** Nested bullet (level 1)
        - *** Deep nested bullet (level 2)
        - ![alt text](path/to/image.png) - untuk add image
        - | Table | Format | - untuk add table
        - --- untuk slide separator

        Args:
            md_content: String markdown content
            output_path: Output file path
        """
        # Split by slide separator
        slides_content = md_content.split("\n---\n")

        for idx, slide_content in enumerate(slides_content):
            if not slide_content.strip():
                continue

            lines = slide_content.split("\n")

            title = None
            content_items = []
            image_paths = []
            table_data = None

            # Detect table
            table_lines = []
            in_table = False
            non_table_lines = []

            for line in lines:
                line_stripped = line.strip()

                # Check if line is part of table
                if line_stripped.startswith("|") and "|" in line_stripped[1:]:
                    in_table = True
                    table_lines.append(line)
                else:
                    if in_table and table_lines:
                        # Table ended, parse it
                        table_text = "\n".join(table_lines)
                        table_data = self.parse_markdown_table(table_text)
                        table_lines = []
                        in_table = False
                    non_table_lines.append(line)

            # Handle table at end of slide
            if table_lines:
                table_text = "\n".join(table_lines)
                table_data = self.parse_markdown_table(table_text)

            # Process non-table lines
            for line in non_table_lines:
                if not line.strip():
                    continue

                original_line = line
                line = line.strip()

                # Header level 1 (# Header) - biasanya untuk slide title utama
                if line.startswith("# ") and not line.startswith("## "):
                    title = line.replace("# ", "").strip()

                # Header level 2 (## Header) - slide title
                elif line.startswith("## ") and not line.startswith("### "):
                    heading = line.replace("## ", "").strip()
                    if not title:
                        title = heading
                    else:
                        content_items.append(
                            {"type": "section", "text": heading, "header_level": 2}
                        )

                # Header level 3 (### Header) - section dalam slide
                elif line.startswith("### ") and not line.startswith("#### "):
                    content_items.append(
                        {
                            "type": "section",
                            "text": line.replace("### ", "").strip(),
                            "header_level": 3,
                        }
                    )

                # Header level 4 (#### Header) - sub-section dalam slide
                elif line.startswith("#### "):
                    content_items.append(
                        {
                            "type": "section",
                            "text": line.replace("#### ", "").strip(),
                            "header_level": 4,
                        }
                    )

                # Bullet point dengan deteksi level (indentasi)
                elif line.startswith("* ") or line.startswith("- "):
                    # Count leading spaces untuk detect nested level
                    leading_spaces = len(original_line) - len(original_line.lstrip())

                    # Detect nested bullet by spaces or asterisks
                    if "**" in line[:5]:
                        level = 2
                        bullet_text = line.lstrip("*- ").strip()
                    elif (
                        line.startswith("  *")
                        or line.startswith("  -")
                        or leading_spaces >= 2
                    ):
                        level = 1
                        bullet_text = line.lstrip("*- ").strip()
                    else:
                        level = 0
                        bullet_text = line.lstrip("*- ").strip()

                    content_items.append(
                        {"type": "bullet", "text": bullet_text, "level": level}
                    )

                # Image reference ![alt](path)
                elif line.startswith("!["):
                    match = re.match(r"!\[.*?\]\((.*?)\)", line)
                    if match:
                        image_paths.append(match.group(1))

                # Regular text
                elif line and not line.startswith("#"):
                    content_items.append({"type": "text", "text": line})

            # Add slide
            content = content_items if content_items else None
            self.add_slide_from_content(
                title=title,
                content=content,
                image_paths=image_paths if image_paths else None,
                table_data=table_data,
                layout_index=1,
            )

        # Save
        self.prs.save(output_path)
        mode = "template-based" if self.use_template else "generated"
        print(f"\n✅ Presentasi berhasil dibuat ({mode}): {output_path}")
        print(f"   Total slides: {len(self.prs.slides)}")

    def convert_to_bytes(self, markdown_content: str) -> bytes:
        """
        Konversi Markdown menjadi PPTX dalam bentuk bytes (in-memory).

        Args:
            markdown_content (str): Konten dalam format Markdown.

        Returns:
            bytes: File PPTX dalam bentuk bytes.
        """
        try:
            # Reset presentation untuk konversi baru
            if self.template_path and Path(self.template_path).exists():
                self.prs = Presentation(self.template_path)
            else:
                self.prs = Presentation()
                self.prs.slide_width = Inches(10)
                self.prs.slide_height = Inches(5.625)

            # Process markdown content (sama seperti convert_markdown tapi tanpa save ke file)
            # Split by slide separator
            slides_content = markdown_content.split("\n---\n")

            for idx, slide_content in enumerate(slides_content):
                if not slide_content.strip():
                    continue

                lines = slide_content.split("\n")

                title = None
                content_items = []
                image_paths = []
                table_data = None

                # Detect table
                table_lines = []
                in_table = False
                non_table_lines = []

                for line in lines:
                    line_stripped = line.strip()

                    # Check if line is part of table
                    if line_stripped.startswith("|") and "|" in line_stripped[1:]:
                        in_table = True
                        table_lines.append(line)
                    else:
                        if in_table and table_lines:
                            # Table ended, parse it
                            table_text = "\n".join(table_lines)
                            table_data = self.parse_markdown_table(table_text)
                            table_lines = []
                            in_table = False
                        non_table_lines.append(line)

                # Handle table at end of slide
                if table_lines:
                    table_text = "\n".join(table_lines)
                    table_data = self.parse_markdown_table(table_text)

                # Process non-table lines
                for line in non_table_lines:
                    if not line.strip():
                        continue

                    original_line = line
                    line = line.strip()

                    # Header level 1 (# Header)
                    if line.startswith("# ") and not line.startswith("## "):
                        title = line.replace("# ", "").strip()

                    # Header level 2 (## Header)
                    elif line.startswith("## ") and not line.startswith("### "):
                        heading = line.replace("## ", "").strip()
                        if not title:
                            title = heading
                        else:
                            content_items.append(
                                {"type": "section", "text": heading, "header_level": 2}
                            )

                    # Header level 3 (### Header)
                    elif line.startswith("### ") and not line.startswith("#### "):
                        content_items.append(
                            {
                                "type": "section",
                                "text": line.replace("### ", "").strip(),
                                "header_level": 3,
                            }
                        )

                    # Header level 4 (#### Header)
                    elif line.startswith("#### "):
                        content_items.append(
                            {
                                "type": "section",
                                "text": line.replace("#### ", "").strip(),
                                "header_level": 4,
                            }
                        )

                    # Bullet point dengan deteksi level
                    elif line.startswith("* ") or line.startswith("- "):
                        leading_spaces = len(original_line) - len(
                            original_line.lstrip()
                        )

                        if "**" in line[:5]:
                            level = 2
                            bullet_text = line.lstrip("*- ").strip()
                        elif (
                            line.startswith("  *")
                            or line.startswith("  -")
                            or leading_spaces >= 2
                        ):
                            level = 1
                            bullet_text = line.lstrip("*- ").strip()
                        else:
                            level = 0
                            bullet_text = line.lstrip("*- ").strip()

                        content_items.append(
                            {"type": "bullet", "text": bullet_text, "level": level}
                        )

                    # Image reference
                    elif line.startswith("!["):
                        match = re.match(r"!\[.*?\]\((.*?)\)", line)
                        if match:
                            image_paths.append(match.group(1))

                    # Regular text
                    elif line and not line.startswith("#"):
                        content_items.append({"type": "text", "text": line})

                # Add slide
                content = content_items if content_items else None
                self.add_slide_from_content(
                    title=title,
                    content=content,
                    image_paths=image_paths if image_paths else None,
                    table_data=table_data,
                    layout_index=1,
                )

            # Save to BytesIO
            bytes_io = BytesIO()
            self.prs.save(bytes_io)
            bytes_io.seek(0)

            print(f"\n✅ Presentasi berhasil dikonversi ke bytes")
            print(f"   Total slides: {len(self.prs.slides)}")

            return bytes_io.getvalue()

        except Exception as e:
            print(f"✗ Error saat konversi ke bytes: {e}")
            return b""

    def get_bytes(self) -> bytes:
        """
        Get current presentation sebagai bytes (tanpa konversi markdown).
        Untuk backward compatibility.

        Returns:
            bytes: File PPTX dalam bentuk bytes.
        """
        try:
            bytes_io = BytesIO()
            self.prs.save(bytes_io)
            bytes_io.seek(0)
            return bytes_io.getvalue()
        except Exception as e:
            print(f"✗ Error saat get bytes: {e}")
            return b""

    def save(self, output_path):
        """Save presentation"""
        self.prs.save(output_path)
        print(f"✅ Saved: {output_path}")
